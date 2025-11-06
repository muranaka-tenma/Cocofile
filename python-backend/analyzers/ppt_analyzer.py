"""
PowerPoint分析モジュール

python-pptxを使用してPowerPointファイルからテキストを抽出します。
"""

import os
from typing import Dict, Any

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


def analyze_ppt(file_path: str) -> Dict[str, Any]:
    """
    PowerPointファイルを分析してテキストを抽出

    Args:
        file_path: PowerPointファイルのパス

    Returns:
        {
            "text": "抽出されたテキスト",
            "slide_count": スライド数,
            "file_size": ファイルサイズ(bytes)
        }
    """
    if not PPTX_AVAILABLE:
        return {
            "error": "python-pptx is not installed",
            "text": "",
            "slide_count": 0,
            "file_size": 0
        }

    if not os.path.exists(file_path):
        return {
            "error": f"File not found: {file_path}",
            "text": "",
            "slide_count": 0,
            "file_size": 0
        }

    try:
        file_size = os.path.getsize(file_path)

        # プレゼンテーション読み込み
        prs = Presentation(file_path)
        slide_count = len(prs.slides)

        extracted_text = []

        # 全スライドからテキスト抽出
        for i, slide in enumerate(prs.slides, start=1):
            extracted_text.append(f"[Slide {i}]")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    extracted_text.append(shape.text)

        full_text = "\n".join(extracted_text)

        return {
            "text": full_text,
            "slide_count": slide_count,
            "file_size": file_size,
            "success": True
        }

    except Exception as e:
        return {
            "error": f"Failed to analyze PowerPoint: {str(e)}",
            "text": "",
            "slide_count": 0,
            "file_size": os.path.getsize(file_path) if os.path.exists(file_path) else 0
        }


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        result = analyze_ppt(sys.argv[1])
        print(result)
    else:
        print("Usage: python ppt_analyzer.py <ppt_file_path>")
